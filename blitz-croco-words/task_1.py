import os
from pathlib import Path
from zipfile import ZipFile
from pptx import Presentation


def getting_words(a: []):
    words = ""
    x = Presentation(a)
    for slide in x.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text: str = shape.text
                if ' ' in text or ':' in text or 'БЛИЦ-КРОКОДИЛ' in text:
                    continue
                words += text + '\n'
    return words

def parsing(name):
    words: list = []
    words_not_split = ""

    with ZipFile(name) as archive:
        for f in archive.namelist():
            words_not_split += (getting_words(archive.open(f)))

    words = words_not_split.split('\n')
    words.pop()

    for i in range(len(words)):
        print(f"[{i + 1}] {words[i]}")


if __name__ == '__main__':
    parsing("croco-blitz-source.zip")
