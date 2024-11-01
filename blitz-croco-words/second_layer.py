from pptx import Presentation

def getting_words(a: []):
    words = ""
    x = Presentation(a)
    for slide in x.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text: str = shape.text
                if ' ' in text or ':' in text or 'БЛИЦ-КРОКОДИЛ' in text:
                    continue
                words += text + '\n'
    return words