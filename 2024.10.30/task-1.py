from pptx import Presentation

prs = Presentation("Zimnyaya_igra_1.pptx")
words = []

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            text: str = shape.text
            if ' ' in text or ':' in text or 'БЛИЦ-КРОКОДИЛ' in text:
                continue
            words.append(text)


with open("words.txt", "w", encoding="UTF-8") as file:
    for i in range(len(words)):
        file.write(str(words[i]) + '\n')
